import os
from typing import List
from pypdf import PdfReader
from PIL import Image

# optional libs
try:
    import docx  # python-docx
    DOCX_AVAILABLE = True
except Exception:
    DOCX_AVAILABLE = False

try:
    import docx2txt
    DOCX2TXT_AVAILABLE = True
except Exception:
    DOCX2TXT_AVAILABLE = False

try:
    import textract  # fallback for older .doc files
    TEXTRACT_AVAILABLE = True
except Exception:
    TEXTRACT_AVAILABLE = False

try:
    from pptx import Presentation  # python-pptx
    PPTX_AVAILABLE = True
except Exception:
    PPTX_AVAILABLE = False

# OCR backends
try:
    import pytesseract
    TESSERACT_AVAILABLE = True
except Exception:
    TESSERACT_AVAILABLE = False

# pdf2image for scanned pdf -> images
try:
    from pdf2image import convert_from_path
    PDF2IMAGE_AVAILABLE = True
except Exception:
    PDF2IMAGE_AVAILABLE = False

# easyocr fallback
EASYOCR_AVAILABLE = False
try:
    import easyocr
    _easyocr_reader = easyocr.Reader(["en"], gpu=False)
    EASYOCR_AVAILABLE = True
except Exception:
    EASYOCR_AVAILABLE = False


def load_txt(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def load_docx(path: str) -> str:
    if DOCX2TXT_AVAILABLE:
        return docx2txt.process(path)
    if DOCX_AVAILABLE:
        doc = docx.Document(path)
        return "\\n".join([p.text for p in doc.paragraphs if p.text])
    raise RuntimeError("python-docx or docx2txt is required to read .docx files. Install python-docx or docx2txt.")


def load_doc(path: str) -> str:
    if TEXTRACT_AVAILABLE:
        txt = textract.process(path)
        if isinstance(txt, bytes):
            try:
                return txt.decode("utf-8", errors="ignore")
            except Exception:
                return str(txt)
        return str(txt)
    else:
        raise RuntimeError("Reading .doc requires 'textract' and some system binaries. Convert .doc to .docx if possible.")


def load_pptx(path: str) -> str:
    if not PPTX_AVAILABLE:
        raise RuntimeError("python-pptx is required to read .pptx. Install python-pptx.")
    pres = Presentation(path)
    texts = []
    for slide in pres.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                if shape.text:
                    texts.append(shape.text)
    return "\\n".join(texts)


def ocr_image_with_tesseract(img: Image.Image) -> str:
    return pytesseract.image_to_string(img)


def ocr_image_with_easyocr(path: str) -> str:
    return "\\n".join(_easyocr_reader.readtext(path, detail=0))


def load_image(path: str) -> str:
    img = Image.open(path).convert("RGB")
    if TESSERACT_AVAILABLE:
        try:
            return ocr_image_with_tesseract(img)
        except Exception:
            pass
    if EASYOCR_AVAILABLE:
        return ocr_image_with_easyocr(path)
    raise RuntimeError("No OCR backend available. Install Tesseract or easyocr.")


def load_pdf(path: str) -> str:
    reader = PdfReader(path)
    pages = []
    for p in reader.pages:
        t = p.extract_text()
        if t:
            pages.append(t)
    joined = "\\n".join(pages).strip()
    if joined:
        return joined

    # fallback: scanned PDF -> convert pages to images & OCR
    if PDF2IMAGE_AVAILABLE:
        images = convert_from_path(path)
        texts = []
        for img in images:
            if TESSERACT_AVAILABLE:
                texts.append(ocr_image_with_tesseract(img))
            elif EASYOCR_AVAILABLE:
                tmp = path + ".tmp_page.png"
                img.save(tmp)
                texts.append(ocr_image_with_easyocr(tmp))
                try:
                    os.remove(tmp)
                except Exception:
                    pass
        return "\\n".join(texts)

    raise RuntimeError(
        "PDF has no extractable text and pdf2image/poppler is not available for OCR. "
        "Install pdf2image/poppler or provide a text-based PDF."
    )


def load_document(file_path: str) -> str:
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".txt":
        return load_txt(file_path)
    if ext == ".pdf":
        return load_pdf(file_path)
    if ext == ".docx":
        return load_docx(file_path)
    if ext == ".doc":
        return load_doc(file_path)
    if ext == ".pptx":
        return load_pptx(file_path)
    if ext in [".png", ".jpg", ".jpeg", ".tiff", ".bmp", ".webp"]:
        return load_image(file_path)
    raise ValueError("Unsupported file type. Supported: TXT, PDF, DOCX, DOC, PPTX, PNG, JPG, JPEG, TIFF, BMP, WEBP.")
